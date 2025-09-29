# translate_ppt_deepl_app.py
#
# Application Streamlit pour traduire des fichiers PowerPoint (.pptx ou .ppt) et Jupyter Notebooks (.ipynb) 
# du français vers l'anglais en utilisant l'API DeepL, tout en préservant la mise en forme.
#
# ⚙️ Prérequis
# - Python 3.9+
# - Packages : streamlit, python-pptx, requests, nbformat
#   pip install streamlit python-pptx requests nbformat
# - Clé API DeepL via la variable d'environnement DEEPL_API_KEY
#   export DEEPL_API_KEY="votre_clef_deepl"
# - (Optionnel pour .ppt) LibreOffice installé avec la commande `soffice` disponible dans le PATH
#
# ▶️ Lancer :
#   streamlit run translate_ppt_deepl_app.py
#
# 📝 Notes
# - PowerPoint : La traduction se fait run par run (éléments de texte formatés) pour préserver la mise en forme.
# - PowerPoint : Les tableaux et les notes des diapositives sont pris en charge.
# - PowerPoint : Les graphiques/SmartArt/objets intégrés ne sont pas modifiables via python-pptx et ne seront pas traduits.
# - PowerPoint : Les fichiers .ppt (ancien format) sont automatiquement convertis en .pptx via LibreOffice si disponible.
# - Jupyter : Seuls les cellules markdown et les commentaires dans le code sont traduits, le code reste intact.

import io
import json
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import List, Tuple

import nbformat
import requests
import streamlit as st
from pptx import Presentation
from pptx.shapes.group import GroupShape

# Charger les variables d'environnement depuis .env si le fichier existe
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    # python-dotenv n'est pas installé, continuer sans
    pass

DEEPL_API_KEY = os.getenv("DEEPL_API_KEY", "")
DEEPL_API_URL = os.getenv("DEEPL_API_URL")  # Permet de forcer l'URL si besoin

# Détection d'URL selon type de clé : les clés Free contiennent "-free" ; sinon, endpoint payant.
if not DEEPL_API_URL:
    if DEEPL_API_KEY.endswith(":fx") or "-free" in DEEPL_API_KEY:
        DEEPL_API_URL = "https://api-free.deepl.com/v2/translate"
    else:
        DEEPL_API_URL = "https://api.deepl.com/v2/translate"

BATCH_SIZE = 45  # par sécurité, rester < 50 textes par requête

# Fonction de traduction par lots via DeepL
def deepl_translate_batch(texts: List[str], source_lang: str = "FR", target_lang: str = "EN-US") -> List[str]:
    if not texts:
        return []
    # Construire un payload 'application/x-www-form-urlencoded' avec répétition de la clé 'text'
    data: List[Tuple[str, str]] = [
        ("auth_key", DEEPL_API_KEY),
        ("source_lang", source_lang),
        ("target_lang", target_lang),
        ("preserve_formatting", "1"),  # aide à garder la casse/ponctuation
    ]
    data.extend(("text", t) for t in texts)

    # Requêtes + gestion d'erreurs simples
    try:
        resp = requests.post(DEEPL_API_URL, data=data, timeout=60)
        resp.raise_for_status()
        js = resp.json()
        translations = js.get("translations", [])
        return [t.get("text", "") for t in translations]
    except requests.HTTPError as e:
        st.error(f"Erreur DeepL ({e.response.status_code}) : {e.response.text}")
        raise
    except Exception as e:
        st.error(f"Erreur de connexion à l'API DeepL : {e}")
        raise


def process_powerpoint_file(in_path: Path, tmpdir: Path, tgt_variant: str, include_notes: bool):
    """Traite un fichier PowerPoint pour la traduction."""

    # Convertir .ppt → .pptx via LibreOffice si nécessaire
    def convert_ppt_to_pptx(ppt_path: Path) -> Path:
        out_dir = ppt_path.parent
        try:
            result = subprocess.run([
                "soffice", "--headless", "--convert-to", "pptx", "--outdir", str(out_dir), str(ppt_path)
            ], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=120)
            if result.returncode != 0:
                raise RuntimeError(result.stderr or result.stdout)
            conv = ppt_path.with_suffix(".pptx")
            if not conv.exists():
                # Certains LibreOffice sortent un nom sans respecter la casse
                candidates = list(out_dir.glob(ppt_path.stem + "*.pptx"))
                if candidates:
                    return candidates[0]
                raise FileNotFoundError("Conversion échouée : fichier .pptx introuvable")
            return conv
        except Exception as e:
            raise RuntimeError(f"Conversion .ppt→.pptx impossible : {e}")

    if in_path.suffix.lower() == ".ppt":
        try:
            st.info("Conversion du .ppt vers .pptx via LibreOffice…")
            in_path = convert_ppt_to_pptx(in_path)
            st.success("Conversion .ppt → .pptx réalisée.")
        except Exception as e:
            st.error(str(e))
            st.stop()

    # Charger la présentation
    prs = Presentation(str(in_path))

    # Utilitaires : itérer sur tous les conteneurs de texte (y compris GroupShapes et Tables)
    def iter_text_frames(shape):
        """Génère tous les text_frames pour un shape (y compris récursif pour groupes et cellules de tableau)."""
        if isinstance(shape, GroupShape):
            for shp in shape.shapes:
                yield from iter_text_frames(shp)
            return

        # Table - utiliser has_table pour une vérification sûre
        if hasattr(shape, 'has_table') and shape.has_table:
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame is not None:
                            yield cell.text_frame
                return
            except (ValueError, AttributeError):
                # Ignorer les erreurs pour les tableaux problématiques
                pass

        # Formes avec texte (auto-shapes, placeholders, text boxes)
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            yield shape.text_frame

    # Collecter tous les runs de texte pour traduction
    RunRef = Tuple  # alias pour lisibilité: (run_obj, original_text)
    run_refs: List[Tuple[object, str]] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            for tf in iter_text_frames(shape):
                for para in tf.paragraphs:
                    for run in para.runs:
                        txt = run.text
                        if txt is not None and txt.strip() != "":
                            run_refs.append((run, txt))

        # Notes de la diapositive
        if include_notes:
            if slide.has_notes_slide:
                notes = slide.notes_slide
                for shape in notes.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                txt = run.text
                                if txt is not None and txt.strip() != "":
                                    run_refs.append((run, txt))

    total_runs = len(run_refs)
    if total_runs == 0:
        st.warning("Aucun texte détecté à traduire.")
        st.stop()

    st.write(f"Segments à traduire : **{total_runs}** (en conservant la mise en forme)")
    
    # Appliquer la traduction par lots en conservant la mise en forme (remplacement run par run)
    progress = st.progress(0, text="Traduction en cours…")
    for start in range(0, total_runs, BATCH_SIZE):
        batch_refs = run_refs[start:start + BATCH_SIZE]
        batch_texts = [txt for (_run, txt) in batch_refs]
        translated = deepl_translate_batch(batch_texts, source_lang="FR", target_lang=tgt_variant)
        for (run, _orig), new_txt in zip(batch_refs, translated):
            run.text = new_txt
        progress.progress(min(1.0, (start + len(batch_refs)) / total_runs), text=f"{min(start + len(batch_refs), total_runs)}/{total_runs} segments")

    # Enregistrer la présentation traduite
    out_name = Path(in_path.name).with_suffix("")
    out_file = tmpdir / f"{out_name}_EN.pptx"
    prs.save(str(out_file))

    st.success("Traduction terminée. Téléchargez votre fichier ci-dessous.")
    with open(out_file, "rb") as f:
        st.download_button(
            label="⬇️ Télécharger le PowerPoint traduit",
            data=f.read(),
            file_name=out_file.name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    # Proposer un bouton pour réinitialiser la session (facilite une autre traduction)
    if st.button("🔁 Traduire un autre fichier"):
        st.experimental_rerun()


def process_jupyter_notebook(in_path: Path, tmpdir: Path, tgt_variant: str):
    """Traite un fichier Jupyter Notebook pour la traduction."""
    
    # Charger le notebook
    try:
        with open(in_path, 'r', encoding='utf-8') as f:
            notebook = nbformat.read(f, as_version=4)
    except Exception as e:
        st.error(f"Erreur lors du chargement du notebook : {e}")
        st.stop()

    # Collecter tous les textes à traduire
    texts_to_translate = []
    text_refs = []  # Références pour remplacer les textes traduits
    
    def extract_comments_from_code(code_text: str) -> List[Tuple[str, int, int]]:
        """Extrait les commentaires d'un code Python et retourne (commentaire, start, end)."""
        comments = []
        lines = code_text.split('\n')
        for i, line in enumerate(lines):
            # Commentaires de ligne (#)
            if '#' in line:
                # Trouver le début du commentaire
                comment_start = line.find('#')
                comment_text = line[comment_start + 1:].strip()
                if comment_text:
                    # Calculer les positions dans le texte original
                    start_pos = sum(len(lines[j]) + 1 for j in range(i)) + comment_start + 1
                    end_pos = start_pos + len(comment_text)
                    comments.append((comment_text, start_pos, end_pos))
        return comments

    # Parcourir toutes les cellules
    for cell_idx, cell in enumerate(notebook.cells):
        if cell.cell_type == 'markdown':
            # Cellules markdown - traduire tout le contenu
            if cell.source and cell.source.strip():
                texts_to_translate.append(cell.source)
                text_refs.append(('markdown', cell_idx, None))
        
        elif cell.cell_type == 'code':
            # Cellules code - extraire seulement les commentaires
            if cell.source:
                comments = extract_comments_from_code(cell.source)
                for comment_text, start_pos, end_pos in comments:
                    if comment_text.strip():
                        texts_to_translate.append(comment_text)
                        text_refs.append(('comment', cell_idx, (start_pos, end_pos)))

    total_texts = len(texts_to_translate)
    if total_texts == 0:
        st.warning("Aucun texte détecté à traduire (markdown ou commentaires).")
        st.stop()

    st.write(f"Segments à traduire : **{total_texts}** (markdown et commentaires)")
    
    # Traduction par lots
    progress = st.progress(0, text="Traduction en cours…")
    translated_texts = []
    
    for start in range(0, total_texts, BATCH_SIZE):
        batch_texts = texts_to_translate[start:start + BATCH_SIZE]
        batch_translated = deepl_translate_batch(batch_texts, source_lang="FR", target_lang=tgt_variant)
        translated_texts.extend(batch_translated)
        progress.progress(min(1.0, (start + len(batch_texts)) / total_texts), 
                         text=f"{min(start + len(batch_texts), total_texts)}/{total_texts} segments")

    # Appliquer les traductions au notebook
    for (text_type, cell_idx, position), translated_text in zip(text_refs, translated_texts):
        if text_type == 'markdown':
            # Remplacer le contenu markdown
            notebook.cells[cell_idx].source = translated_text
        elif text_type == 'comment':
            # Remplacer le commentaire dans le code
            cell = notebook.cells[cell_idx]
            original_code = cell.source
            start_pos, end_pos = position
            
            # Remplacer le commentaire dans le code original
            new_code = original_code[:start_pos] + translated_text + original_code[end_pos:]
            cell.source = new_code

    # Enregistrer le notebook traduit
    out_name = Path(in_path.name).with_suffix("")
    out_file = tmpdir / f"{out_name}_EN.ipynb"
    
    try:
        with open(out_file, 'w', encoding='utf-8') as f:
            nbformat.write(notebook, f)
    except Exception as e:
        st.error(f"Erreur lors de la sauvegarde du notebook : {e}")
        st.stop()

    st.success("Traduction terminée. Téléchargez votre fichier ci-dessous.")
    with open(out_file, "rb") as f:
        st.download_button(
            label="⬇️ Télécharger le Notebook traduit",
            data=f.read(),
            file_name=out_file.name,
            mime="application/x-ipynb+json",
        )

    # Proposer un bouton pour réinitialiser la session
    if st.button("🔁 Traduire un autre fichier"):
        st.experimental_rerun()

st.set_page_config(page_title="Traduire FR → EN (DeepL)", page_icon="🌐", layout="centered")

# Navigation par pages
page = st.sidebar.selectbox("Choisir le type de fichier à traduire", ["PowerPoint", "Jupyter Notebook"])

if page == "PowerPoint":
    st.title("🗂️ Traduire un PowerPoint FR → EN (DeepL)")
    st.write("Téléversez un fichier **.pptx** ou **.ppt**. La mise en forme est préservée ; seul le texte est traduit en anglais.")
elif page == "Jupyter Notebook":
    st.title("📓 Traduire un Jupyter Notebook FR → EN (DeepL)")
    st.write("Téléversez un fichier **.ipynb**. Seuls les cellules markdown et les commentaires dans le code sont traduits ; le code reste intact.")

with st.expander("Paramètres avancés"):
    tgt_variant = st.selectbox("Variante d'anglais", ["EN-US", "EN-GB"], index=0)
    if page == "PowerPoint":
        include_notes = st.checkbox("Traduire les notes des diapositives", value=True)

if page == "PowerPoint":
    uploaded = st.file_uploader("Choisir un fichier PowerPoint", type=["pptx", "ppt"])
elif page == "Jupyter Notebook":
    uploaded = st.file_uploader("Choisir un fichier Jupyter Notebook", type=["ipynb"]) 

if uploaded is not None:
    if not DEEPL_API_KEY:
        st.error("La variable d'environnement **DEEPL_API_KEY** n'est pas définie.")
        st.stop()

    # Sauvegarder le fichier uploadé dans un dossier temporaire
    tmpdir = Path(tempfile.mkdtemp(prefix="translate-"))
    in_path = tmpdir / uploaded.name
    with open(in_path, "wb") as f:
        f.write(uploaded.getbuffer())

    if page == "PowerPoint":
        # Logique PowerPoint
        process_powerpoint_file(in_path, tmpdir, tgt_variant, include_notes)
    elif page == "Jupyter Notebook":
        # Logique Jupyter Notebook
        process_jupyter_notebook(in_path, tmpdir, tgt_variant)

